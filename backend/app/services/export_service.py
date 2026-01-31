from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
from datetime import datetime


class RoadmapExporter:
    """Service for exporting roadmap to PowerPoint"""

    def __init__(self, categories):
        self.categories = categories

        # Color mapping from Tailwind to PowerPoint RGB
        # Maps to the exact Tailwind -700 colors used in the frontend
        # Frontend uses: bg-{color}-900/30 border-{color}-700/50 text-{color}-400
        self.color_map = {
            "blue": RGBColor(29, 78, 216),  # blue-700
            "green": RGBColor(21, 128, 61),  # green-700
            "purple": RGBColor(126, 34, 206),  # purple-700
            "yellow": RGBColor(161, 98, 7),  # yellow-700
            "pink": RGBColor(190, 24, 93),  # pink-700
            "indigo": RGBColor(67, 56, 202),  # indigo-700
            "red": RGBColor(185, 28, 28),  # red-700
            "orange": RGBColor(194, 65, 12),  # orange-700
            "teal": RGBColor(15, 118, 110),  # teal-700
            "cyan": RGBColor(14, 116, 144),  # cyan-700
        }

        # Text colors (Tailwind -400 variants)
        self.text_color_map = {
            "blue": RGBColor(96, 165, 250),  # blue-400
            "green": RGBColor(74, 222, 128),  # green-400
            "purple": RGBColor(192, 132, 252),  # purple-400
            "yellow": RGBColor(250, 204, 21),  # yellow-400
            "pink": RGBColor(244, 114, 182),  # pink-400
            "indigo": RGBColor(129, 140, 248),  # indigo-400
            "red": RGBColor(248, 113, 113),  # red-400
            "orange": RGBColor(251, 146, 60),  # orange-400
            "teal": RGBColor(45, 212, 191),  # teal-400
            "cyan": RGBColor(34, 211, 238),  # cyan-400
        }

        self.color_names = [
            "blue",
            "green",
            "purple",
            "yellow",
            "pink",
            "indigo",
            "red",
            "orange",
            "teal",
            "cyan",
        ]

        # Priority colors
        self.priority_colors = {
            "High": RGBColor(220, 38, 38),  # Red
            "Medium": RGBColor(234, 179, 8),  # Yellow
            "Low": RGBColor(34, 197, 94),  # Green
        }

    def generate_pptx(self):
        """Generate PowerPoint presentation and return as BytesIO buffer"""
        # Create presentation with 16:9 widescreen (50% larger)
        prs = Presentation()
        prs.slide_width = Inches(20)
        prs.slide_height = Inches(11.25)

        # Get features with dates
        features_with_dates = self._get_features_with_dates()

        if not features_with_dates:
            raise ValueError("No features with release dates found")

        # Generate quarter range (changed from monthly to quarterly)
        quarters = self._generate_quarter_range(features_with_dates)

        # Create title slide
        self._create_title_slide(prs, quarters)

        # Split timeline into 4-quarter chunks and create slides
        chunks = list(self._chunk_timeline(quarters, chunk_size=4))
        for i, quarters_chunk in enumerate(chunks):
            self._create_timeline_slide(
                prs, quarters_chunk, i, len(chunks), features_with_dates
            )

        # Save to BytesIO buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

    def _create_title_slide(self, prs, quarters):
        """Create title slide with metadata"""
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        # Use blank slide layout to have full control
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background to match web app (Tailwind gray-800: RGB(31, 41, 55))
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 41, 55)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(18), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Product Roadmap"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(243, 244, 246)  # gray-100
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.8), Inches(18), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame

        # Format quarters for subtitle
        start_quarter = self._format_quarter_long(quarters[0])
        end_quarter = self._format_quarter_long(quarters[-1])

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"{start_quarter} - {end_quarter}"
        subtitle_para.font.size = Pt(27)
        subtitle_para.font.color.rgb = RGBColor(209, 213, 219)  # gray-300
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Generated date
        date_para = subtitle_frame.add_paragraph()
        date_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        date_para.font.size = Pt(21)
        date_para.font.color.rgb = RGBColor(209, 213, 219)  # gray-300
        date_para.alignment = PP_ALIGN.CENTER

        # Add legend at bottom (50% larger)
        legend_top = Inches(7.5)
        legend_left = Inches(5)

        # Legend title
        legend_title_box = slide.shapes.add_textbox(
            legend_left, legend_top, Inches(10), Inches(0.6)
        )
        legend_title_frame = legend_title_box.text_frame
        legend_title_para = legend_title_frame.paragraphs[0]
        legend_title_para.text = "Legend:"
        legend_title_para.font.size = Pt(24)
        legend_title_para.font.bold = True
        legend_title_para.font.color.rgb = RGBColor(200, 200, 200)

        # Signed off example (solid border)
        signoff_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            legend_left + Inches(0.3),
            legend_top + Inches(0.75),
            Inches(2.4),
            Inches(0.6),
        )
        signoff_shape.fill.solid()
        signoff_shape.fill.fore_color.rgb = RGBColor(29, 78, 216)  # blue-700
        signoff_shape.fill.transparency = 0.5
        signoff_shape.line.color.rgb = RGBColor(29, 78, 216)
        signoff_shape.line.width = Pt(2.25)

        signoff_text = signoff_shape.text_frame
        signoff_para = signoff_text.paragraphs[0]
        signoff_para.text = "Signed Off"
        signoff_para.font.size = Pt(18)
        signoff_para.font.color.rgb = RGBColor(96, 165, 250)  # blue-400
        signoff_para.alignment = PP_ALIGN.CENTER

        # Pending estimation example (dashed border with hourglass)
        pending_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            legend_left + Inches(5.5),
            legend_top + Inches(0.75),
            Inches(3.3),
            Inches(0.6),
        )
        pending_shape.fill.solid()
        pending_shape.fill.fore_color.rgb = RGBColor(126, 34, 206)  # purple-700
        pending_shape.fill.transparency = 0.5
        pending_shape.line.color.rgb = RGBColor(126, 34, 206)
        pending_shape.line.width = Pt(3)
        pending_shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        pending_text = pending_shape.text_frame
        pending_para = pending_text.paragraphs[0]
        pending_para.text = "⏳ Pending Estimation"
        pending_para.font.size = Pt(18)
        pending_para.font.color.rgb = RGBColor(192, 132, 252)  # purple-400
        pending_para.alignment = PP_ALIGN.CENTER

    def _create_timeline_slide(
        self, prs, quarters_chunk, chunk_index, total_chunks, features_with_dates
    ):
        """Create a timeline slide for a chunk of quarters"""
        # Use blank slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background to match web app (Tailwind gray-800: RGB(31, 41, 55))
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 41, 55)

        # Add title (all dimensions 50% larger)
        title_left = Inches(0.75)
        title_top = Inches(0.45)
        title_width = Inches(18.5)
        title_height = Inches(0.9)

        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]

        start_quarter = self._format_quarter_long(quarters_chunk[0])
        end_quarter = self._format_quarter_long(quarters_chunk[-1])
        title_text = f"Product Roadmap: {start_quarter} - {end_quarter}"

        if total_chunks > 1:
            title_text += f" (Slide {chunk_index + 1} of {total_chunks})"

        title_para.text = title_text
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Calculate dimensions (all 50% larger)
        category_col_width = Inches(2.7)
        timeline_start_left = Inches(0.75) + category_col_width
        timeline_width = Inches(18.5) - category_col_width
        quarter_width = timeline_width / len(quarters_chunk)

        header_top = Inches(1.65)
        header_height = Inches(0.6)

        # Draw quarter headers
        for i, quarter in enumerate(quarters_chunk):
            quarter_left = timeline_start_left + (i * quarter_width)

            quarter_box = slide.shapes.add_textbox(
                quarter_left, header_top, quarter_width, header_height
            )
            quarter_frame = quarter_box.text_frame
            quarter_para = quarter_frame.paragraphs[0]
            quarter_para.text = self._format_quarter_short(quarter)
            quarter_para.font.size = Pt(21)
            quarter_para.font.bold = True
            quarter_para.font.color.rgb = RGBColor(200, 200, 200)
            quarter_para.alignment = PP_ALIGN.CENTER

        # Get categories with features in this timeline
        categories_with_features = []
        for cat in self.categories:
            cat_features = [
                f
                for f in features_with_dates
                if f["category_id"] == cat.id and f["release_quarter"] in quarters_chunk
            ]
            if cat_features or any(
                f["category_id"] == cat.id for f in features_with_dates
            ):
                categories_with_features.append(cat)

        # Draw category rows and features (all 50% larger)
        row_height = Inches(1.35)
        content_top = header_top + header_height + Inches(0.15)

        for cat_index, category in enumerate(categories_with_features):
            row_top = content_top + (cat_index * row_height)

            # Category name
            cat_box = slide.shapes.add_textbox(
                Inches(0.75), row_top, category_col_width - Inches(0.15), row_height
            )
            cat_frame = cat_box.text_frame
            cat_frame.word_wrap = True
            cat_para = cat_frame.paragraphs[0]
            cat_para.text = category.name
            cat_para.font.size = Pt(18)
            cat_para.font.bold = True
            cat_para.font.color.rgb = RGBColor(200, 200, 200)

            # Get color for this category
            color_name = self.color_names[cat_index % len(self.color_names)]
            fill_color = self.color_map[color_name]
            text_color = self.text_color_map[color_name]

            # Draw features for each quarter
            for quarter_index, quarter in enumerate(quarters_chunk):
                quarter_features = [
                    f
                    for f in features_with_dates
                    if f["category_id"] == category.id and f["release_quarter"] == quarter
                ]

                quarter_left = timeline_start_left + (quarter_index * quarter_width)

                # Stack features vertically if multiple in same quarter
                feature_height = (
                    Inches(0.975) if len(quarter_features) <= 1 else Inches(0.6)
                )
                feature_spacing = Inches(0.075)

                for feat_index, feature in enumerate(
                    quarter_features[:3]
                ):  # Max 3 features per cell
                    feature_top = (
                        row_top
                        + Inches(0.15)
                        + (feat_index * (feature_height + feature_spacing))
                    )

                    self._add_feature_card(
                        slide,
                        feature,
                        quarter_left + Inches(0.075),
                        feature_top,
                        quarter_width - Inches(0.15),
                        feature_height,
                        fill_color,
                        text_color,
                    )

    def _add_feature_card(
        self, slide, feature, left, top, width, height, fill_color, text_color
    ):
        """Add a feature card to the slide"""
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        # Add rounded rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )

        # Set fill color with transparency to match web (bg-{color}-900/30)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.fill.transparency = 0.5  # 50% transparency to match web appearance

        # Check if pending estimation (no engineering signoff)
        is_pending = not feature.get("engineering_signoff", True)

        # Set border color and style
        shape.line.color.rgb = fill_color
        shape.line.transparency = 0.3

        if is_pending:
            # Dashed border for pending estimation (50% thicker)
            shape.line.width = Pt(3)
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        else:
            # Solid border for signed off (50% thicker)
            shape.line.width = Pt(2.25)

        # Add text (all margins 50% larger)
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(6)
        text_frame.margin_right = Pt(6)
        text_frame.margin_top = Pt(4.5)
        text_frame.margin_bottom = Pt(4.5)

        # Clear default paragraph
        text_frame.clear()

        # Title (bold) - use text_color to match web (text-{color}-400) - 50% larger font
        # Add hourglass emoji for pending estimation
        p_title = text_frame.paragraphs[0]
        title_prefix = "⏳ " if is_pending else ""
        title_text = feature["title"][:30] + (
            "..." if len(feature["title"]) > 30 else ""
        )
        p_title.text = title_prefix + title_text
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = text_color

        # Priority only (no ID, no T-shirt sizing) - 50% larger font
        p_meta = text_frame.add_paragraph()
        priority_letter = feature["priority"][0]  # H, M, or L
        p_meta.text = f"{priority_letter}"
        p_meta.font.size = Pt(14)
        p_meta.font.bold = True

        # Color code priority
        priority_color = self.priority_colors.get(
            feature["priority"], RGBColor(200, 200, 200)
        )
        p_meta.font.color.rgb = priority_color

    def _get_features_with_dates(self):
        """Extract all features with release dates and category info, converting to quarters"""
        features = []
        for category in self.categories:
            for feature in category.features:
                if feature.release_date:
                    # Convert YYYY-MM to YYYY-QN
                    release_quarter = self._date_to_quarter(feature.release_date)
                    features.append(
                        {
                            "id": feature.id,
                            "title": feature.title,
                            "priority": feature.priority.value,
                            "release_date": feature.release_date,
                            "release_quarter": release_quarter,
                            "category_id": category.id,
                            "category_name": category.name,
                            "engineering_signoff": feature.engineering_signoff,
                        }
                    )
        return features

    def _generate_quarter_range(self, features):
        """Generate quarter range from earliest to latest + 1 quarter"""
        if not features:
            return []

        quarters = [f["release_quarter"] for f in features]
        quarters_set = sorted(set(quarters))

        if not quarters_set:
            return []

        # Parse start and end quarters
        start_year, start_q = self._parse_quarter(quarters_set[0])
        end_year, end_q = self._parse_quarter(quarters_set[-1])

        # Add 1 quarter to end
        end_q += 1
        if end_q > 4:
            end_q = 1
            end_year += 1

        # Generate all quarters between start and end
        quarters_list = []
        current_year, current_q = start_year, start_q

        while (current_year < end_year) or (
            current_year == end_year and current_q <= end_q
        ):
            quarters_list.append(f"{current_year}-Q{current_q}")
            current_q += 1
            if current_q > 4:
                current_q = 1
                current_year += 1

        return quarters_list

    def _chunk_timeline(self, quarters, chunk_size=4):
        """Split timeline into chunks for multiple slides"""
        for i in range(0, len(quarters), chunk_size):
            yield quarters[i : i + chunk_size]

    def _date_to_quarter(self, date_str):
        """Convert YYYY-MM to YYYY-QN"""
        year, month = date_str.split("-")
        month_num = int(month)
        quarter = (month_num - 1) // 3 + 1  # Math.ceil equivalent
        return f"{year}-Q{quarter}"

    def _parse_quarter(self, quarter_str):
        """Parse YYYY-QN to (year, quarter_num)"""
        year, q = quarter_str.split("-Q")
        return int(year), int(q)

    def _format_quarter_short(self, quarter_str):
        """Format YYYY-QN to 'QN YY' (e.g., 'Q1 24')"""
        year, q = quarter_str.split("-Q")
        year_short = year[-2:]
        return f"Q{q} {year_short}"

    def _format_quarter_long(self, quarter_str):
        """Format YYYY-QN to 'QN YYYY' (e.g., 'Q1 2024')"""
        year, q = quarter_str.split("-Q")
        return f"Q{q} {year}"


class PRDExporter:
    """Service for exporting PRD (all categories and features) to Excel and Word"""

    def __init__(self, categories):
        self.categories = categories

    def generate_excel(self):
        """Generate Excel file with all categories and features"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet

        # Create a sheet for each category (no summary sheet for backup/restore)
        has_sheets = False
        for category in self.categories:
            if category.features:
                sheet = wb.create_sheet(self._sanitize_sheet_name(category.name))
                self._create_category_sheet(sheet, category)
                has_sheets = True

        # If no sheets were created, raise an error
        if not has_sheets:
            raise ValueError("No features found to export. Add features to categories before exporting.")

        # Save to BytesIO buffer
        buffer = BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer

    def _sanitize_sheet_name(self, name):
        """Sanitize sheet name to meet Excel requirements"""
        # Excel sheet names can't be longer than 31 chars and can't contain: \ / ? * [ ]
        invalid_chars = ["\\", "/", "?", "*", "[", "]"]
        sanitized = name
        for char in invalid_chars:
            sanitized = sanitized.replace(char, "")
        return sanitized[:31]

    def _create_summary_sheet(self, sheet):
        """Create summary sheet with overview of all categories"""
        from openpyxl.styles import Font, PatternFill, Alignment

        # Title
        sheet["A1"] = "Product Requirements Document (PRD)"
        sheet["A1"].font = Font(size=16, bold=True)
        sheet["A2"] = f'Generated: {datetime.now().strftime("%B %d, %Y at %I:%M %p")}'
        sheet["A2"].font = Font(size=10, italic=True)

        # Headers
        headers = [
            "Category",
            "Description",
            "Total Features",
            "High Priority",
            "Medium Priority",
            "Low Priority",
        ]
        for col, header in enumerate(headers, start=1):
            cell = sheet.cell(row=4, column=col, value=header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(
                start_color="366092", end_color="366092", fill_type="solid"
            )
            cell.alignment = Alignment(horizontal="center", vertical="center")

        # Data rows
        row = 5
        for category in self.categories:
            features = category.features
            high_count = sum(1 for f in features if f.priority.value == "High")
            medium_count = sum(1 for f in features if f.priority.value == "Medium")
            low_count = sum(1 for f in features if f.priority.value == "Low")

            sheet.cell(row=row, column=1, value=category.name)
            sheet.cell(row=row, column=2, value=category.description)
            sheet.cell(row=row, column=3, value=len(features))
            sheet.cell(row=row, column=4, value=high_count)
            sheet.cell(row=row, column=5, value=medium_count)
            sheet.cell(row=row, column=6, value=low_count)
            row += 1

        # Auto-adjust column widths
        for col in range(1, 7):
            sheet.column_dimensions[chr(64 + col)].width = 20

    def _create_category_sheet(self, sheet, category):
        """Create a sheet for a specific category with all its features"""
        from openpyxl.styles import Font, PatternFill, Alignment

        # Category header
        sheet["A1"] = category.name
        sheet["A1"].font = Font(size=14, bold=True)
        sheet["A2"] = category.description
        sheet["A2"].font = Font(size=10, italic=True)
        sheet.merge_cells("A2:K2")

        # Column headers
        headers = [
            "ID",
            "Title",
            "Priority",
            "Description",
            "KPI",
            "Customer",
            "Eng. Comment",
            "Signoff",
            "Complexity",
            "Release Date",
        ]
        for col, header in enumerate(headers, start=1):
            cell = sheet.cell(row=4, column=col, value=header)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(
                start_color="366092", end_color="366092", fill_type="solid"
            )
            cell.alignment = Alignment(
                horizontal="center", vertical="center", wrap_text=True
            )

        # Feature rows
        row = 5
        for feature in category.features:
            sheet.cell(row=row, column=1, value=feature.id)
            sheet.cell(row=row, column=2, value=feature.title)
            sheet.cell(row=row, column=3, value=feature.priority.value)
            sheet.cell(row=row, column=4, value=feature.description)
            sheet.cell(row=row, column=5, value=feature.kpi)
            sheet.cell(row=row, column=6, value=feature.customer_name)
            sheet.cell(row=row, column=7, value=feature.engineering_comment)
            sheet.cell(
                row=row, column=8, value="Yes" if feature.engineering_signoff else "No"
            )
            sheet.cell(row=row, column=9, value=feature.engineering_complexity.value)

            # Format release date
            if feature.release_date:
                try:
                    year, month = feature.release_date.split("-")
                    release_date_obj = datetime(int(year), int(month), 1)
                    sheet.cell(
                        row=row, column=10, value=release_date_obj.strftime("%B %Y")
                    )
                except (ValueError, TypeError, AttributeError):
                    sheet.cell(row=row, column=10, value=feature.release_date)
            else:
                sheet.cell(row=row, column=10, value="")

            # Apply text wrapping to description and comment columns
            sheet.cell(row=row, column=4).alignment = Alignment(
                wrap_text=True, vertical="top"
            )
            sheet.cell(row=row, column=7).alignment = Alignment(
                wrap_text=True, vertical="top"
            )

            row += 1

        # Auto-adjust column widths
        column_widths = [10, 30, 12, 40, 30, 20, 30, 10, 12, 15]
        for col, width in enumerate(column_widths, start=1):
            sheet.column_dimensions[chr(64 + col)].width = width

    def generate_word(self):
        """Generate Word document with all categories and features"""
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Set narrow margins to fit table within page
        section = doc.sections[0]
        section.left_margin = Inches(0.5)
        section.right_margin = Inches(0.5)
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)

        # Set document font to Montserrat
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Montserrat'
        font.size = Pt(11)

        # Title
        title = doc.add_heading("Product Requirements Document", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].font.name = 'Montserrat'
        title.runs[0].font.size = Pt(24)

        # Metadata
        metadata = doc.add_paragraph(
            f'Generated: {datetime.now().strftime("%B %d, %Y")}'
        )
        metadata.alignment = WD_ALIGN_PARAGRAPH.CENTER
        metadata.runs[0].italic = True
        metadata.runs[0].font.name = 'Montserrat'
        metadata.runs[0].font.size = Pt(10)

        doc.add_paragraph()  # Spacing

        # Table of Contents
        toc_heading = doc.add_heading("Table of Contents", 1)
        toc_heading.runs[0].font.name = 'Montserrat'
        for idx, category in enumerate(self.categories, start=1):
            toc_para = doc.add_paragraph(f"{idx}. {category.name}")
            toc_para.runs[0].font.name = 'Montserrat'
            toc_para.runs[0].font.size = Pt(11)

        doc.add_page_break()

        # Add each category
        for idx, category in enumerate(self.categories):
            self._add_category_to_doc(doc, category, idx + 1)
            if idx < len(self.categories) - 1:  # Don't add page break after last category
                doc.add_page_break()

        # Save to BytesIO buffer
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer

    def _add_category_to_doc(self, doc, category, section_number):
        """Add a category section to the Word document with single table format"""
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        # Category title
        heading = doc.add_heading(f"{section_number}. {category.name}", 1)
        heading.runs[0].font.name = 'Montserrat'
        heading.runs[0].font.size = Pt(18)
        heading.runs[0].bold = True

        # Category description (always add, even if empty)
        desc_text = category.description if category.description else "No description provided."
        desc_para = doc.add_paragraph(desc_text)
        desc_para.runs[0].font.name = 'Montserrat'
        desc_para.runs[0].font.size = Pt(11)
        desc_para.runs[0].font.color.rgb = RGBColor(75, 85, 99)  # gray-600
        doc.add_paragraph()  # Spacing

        # Features
        features = category.features
        if not features:
            no_features_para = doc.add_paragraph("No features in this category.")
            no_features_para.runs[0].font.name = 'Montserrat'
            return

        # Create a single table with features as rows
        # Columns: ID, Title, Priority, Description, KPI, Customer
        num_cols = 6
        table = doc.add_table(rows=1 + len(features), cols=num_cols)
        table.style = 'Light Grid Accent 1'

        # Set table to autofit
        table.autofit = False
        table.allow_autofit = False

        # Set column widths (in inches) - optimized for 0.5" margins (7.5" usable width)
        col_widths = [0.5, 1.5, 0.7, 2.7, 1.1, 0.9]  # Total 7.4 inches
        for idx, width in enumerate(col_widths):
            for cell in table.columns[idx].cells:
                cell.width = Inches(width)

        # Header row
        headers = ["ID", "Title", "Priority", "Description", "KPI", "Customer"]
        header_cells = table.rows[0].cells
        for idx, header_text in enumerate(headers):
            cell = header_cells[idx]
            cell.text = header_text
            # Style header
            paragraph = cell.paragraphs[0]
            run = paragraph.runs[0]
            run.font.bold = True
            run.font.name = 'Montserrat'
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.no_proof = True  # Disable spell check and hyphenation

            # Disable word breaking in header
            paragraph.paragraph_format.widow_control = False

            # Set cell background color (dark blue)
            shading_elm = OxmlElement('w:shd')
            shading_elm.set(qn('w:fill'), '366092')
            cell._element.get_or_add_tcPr().append(shading_elm)

            # Center align header text
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Data rows
        for row_idx, feature in enumerate(features, start=1):
            row_cells = table.rows[row_idx].cells

            # Set cell values (only product-focused columns)
            values = [
                feature.id,
                feature.title,
                feature.priority.value,
                feature.description or "-",
                feature.kpi or "-",
                feature.customer_name or "-",
            ]

            for col_idx, value in enumerate(values):
                cell = row_cells[col_idx]
                cell.text = str(value)

                # Style cell content
                paragraph = cell.paragraphs[0]
                run = paragraph.runs[0]
                run.font.name = 'Montserrat'
                run.font.size = Pt(8)

                # Disable word breaking/hyphenation
                paragraph.paragraph_format.widow_control = False

                # Prevent word breaks by setting language to None (disables hyphenation)
                run.font.no_proof = True

                # Vertical alignment
                cell.vertical_alignment = 1  # Center vertically

                # Priority column: color code and center
                if col_idx == 2:  # Priority column
                    run.font.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    if value == "High":
                        run.font.color.rgb = RGBColor(185, 28, 28)  # red-700
                    elif value == "Medium":
                        run.font.color.rgb = RGBColor(161, 98, 7)  # yellow-700
                    elif value == "Low":
                        run.font.color.rgb = RGBColor(21, 128, 61)  # green-700

        doc.add_paragraph()  # Spacing after table

    def _format_release_date(self, date_str):
        """Format YYYY-MM to readable format"""
        try:
            year, month = date_str.split("-")
            date_obj = datetime(int(year), int(month), 1)
            return date_obj.strftime("%B %Y")
        except (ValueError, TypeError, AttributeError):
            return date_str
