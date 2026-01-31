import pytest
from io import BytesIO
from app.services.export_service import RoadmapExporter, PRDExporter
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document


@pytest.mark.unit
class TestRoadmapExporter:
    """Unit tests for RoadmapExporter service"""

    def test_roadmap_exporter_initialization(self, client, multiple_categories):
        """Test RoadmapExporter initialization"""
        exporter = RoadmapExporter(multiple_categories)

        assert exporter.categories == multiple_categories
        assert len(exporter.color_map) == 10
        assert len(exporter.text_color_map) == 10
        assert len(exporter.color_names) == 10

    def test_get_features_with_dates(self, client, category_with_features):
        """Test extracting features with release dates and quarters"""
        category, features = category_with_features
        exporter = RoadmapExporter([category])

        features_with_dates = exporter._get_features_with_dates()

        assert len(features_with_dates) > 0
        assert all(f['release_date'] is not None for f in features_with_dates)
        assert all('category_name' in f for f in features_with_dates)
        assert all('release_quarter' in f for f in features_with_dates)
        assert all('engineering_signoff' in f for f in features_with_dates)
        # Verify quarter format is YYYY-QN
        assert all(f['release_quarter'].count('-Q') == 1 for f in features_with_dates)

    def test_generate_quarter_range(self, client, category_with_features):
        """Test generating quarter range from features"""
        category, features = category_with_features
        exporter = RoadmapExporter([category])

        features_with_dates = exporter._get_features_with_dates()
        quarters = exporter._generate_quarter_range(features_with_dates)

        assert len(quarters) > 0
        assert all('-Q' in quarter for quarter in quarters)  # Format: YYYY-QN
        # Verify quarters are sequential
        for i in range(len(quarters) - 1):
            year1, q1 = exporter._parse_quarter(quarters[i])
            year2, q2 = exporter._parse_quarter(quarters[i + 1])
            # Next quarter should be +1 quarter or same quarter next year
            assert (year1 == year2 and q2 == q1 + 1) or (year2 == year1 + 1 and q1 == 4 and q2 == 1)

    def test_generate_pptx(self, client, category_with_features):
        """Test generating PowerPoint presentation with quarterly roadmap"""
        category, features = category_with_features
        exporter = RoadmapExporter([category])

        pptx_buffer = exporter.generate_pptx()

        assert isinstance(pptx_buffer, BytesIO)
        pptx_buffer.seek(0)

        # Verify it's a valid presentation
        prs = Presentation(pptx_buffer)
        assert len(prs.slides) >= 1

        # Verify slide dimensions are 50% larger (20" x 11.25" instead of 13.33" x 7.5")
        from pptx.util import Inches
        assert prs.slide_width == Inches(20)
        assert prs.slide_height == Inches(11.25)

    def test_pptx_includes_legend(self, client, category_with_features):
        """Test that PowerPoint includes legend for pending estimation"""
        category, features = category_with_features
        exporter = RoadmapExporter([category])

        pptx_buffer = exporter.generate_pptx()
        pptx_buffer.seek(0)

        prs = Presentation(pptx_buffer)

        # Title slide should be first
        title_slide = prs.slides[0]

        # Verify legend text exists in title slide
        slide_text = []
        for shape in title_slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)

        full_text = ' '.join(slide_text)
        assert 'Legend' in full_text or 'Pending Estimation' in full_text or 'Signed Off' in full_text

    def test_pending_estimation_indicator(self, client, sample_project):
        """Test that features without engineering signoff show pending indicator"""
        from app.models import Category, Feature
        from app import db

        # Create category with features - one signed off, one pending
        category = Category(
            id='test-cat',
            name='Test Category',
            description='Test',
            project_id=sample_project.id
        )
        db.session.add(category)

        signed_off_feature = Feature(
            id='F-SIGNED',
            title='Signed Off Feature',
            category_id='test-cat',
            release_date='2024-06',
            engineering_signoff=True
        )

        pending_feature = Feature(
            id='F-PENDING',
            title='Pending Feature',
            category_id='test-cat',
            release_date='2024-06',
            engineering_signoff=False
        )

        db.session.add(signed_off_feature)
        db.session.add(pending_feature)
        db.session.commit()

        exporter = RoadmapExporter([category])
        features_with_dates = exporter._get_features_with_dates()

        # Find our test features
        signed_off = next(f for f in features_with_dates if f['id'] == 'F-SIGNED')
        pending = next(f for f in features_with_dates if f['id'] == 'F-PENDING')

        assert signed_off['engineering_signoff'] == True
        assert pending['engineering_signoff'] == False

    def test_chunk_timeline(self, client):
        """Test chunking timeline into 4-quarter segments"""
        from app.services.export_service import RoadmapExporter

        quarters = ['2024-Q1', '2024-Q2', '2024-Q3', '2024-Q4',
                    '2025-Q1', '2025-Q2', '2025-Q3', '2025-Q4']

        exporter = RoadmapExporter([])
        chunks = list(exporter._chunk_timeline(quarters, chunk_size=4))

        assert len(chunks) == 2
        assert len(chunks[0]) == 4
        assert len(chunks[1]) == 4

    def test_date_to_quarter_conversion(self, client):
        """Test converting YYYY-MM dates to YYYY-QN quarters"""
        from app.services.export_service import RoadmapExporter

        exporter = RoadmapExporter([])

        # Test Q1 (Jan, Feb, Mar)
        assert exporter._date_to_quarter('2024-01') == '2024-Q1'
        assert exporter._date_to_quarter('2024-02') == '2024-Q1'
        assert exporter._date_to_quarter('2024-03') == '2024-Q1'

        # Test Q2 (Apr, May, Jun)
        assert exporter._date_to_quarter('2024-04') == '2024-Q2'
        assert exporter._date_to_quarter('2024-05') == '2024-Q2'
        assert exporter._date_to_quarter('2024-06') == '2024-Q2'

        # Test Q3 (Jul, Aug, Sep)
        assert exporter._date_to_quarter('2024-07') == '2024-Q3'
        assert exporter._date_to_quarter('2024-08') == '2024-Q3'
        assert exporter._date_to_quarter('2024-09') == '2024-Q3'

        # Test Q4 (Oct, Nov, Dec)
        assert exporter._date_to_quarter('2024-10') == '2024-Q4'
        assert exporter._date_to_quarter('2024-11') == '2024-Q4'
        assert exporter._date_to_quarter('2024-12') == '2024-Q4'

    def test_format_quarter(self, client):
        """Test quarter formatting functions"""
        from app.services.export_service import RoadmapExporter

        exporter = RoadmapExporter([])

        # Test short format
        assert exporter._format_quarter_short('2024-Q1') == 'Q1 24'
        assert exporter._format_quarter_short('2025-Q3') == 'Q3 25'

        # Test long format
        assert exporter._format_quarter_long('2024-Q1') == 'Q1 2024'
        assert exporter._format_quarter_long('2025-Q3') == 'Q3 2025'


@pytest.mark.unit
class TestPRDExporter:
    """Unit tests for PRDExporter service"""

    def test_prd_exporter_initialization(self, client, multiple_categories):
        """Test PRDExporter initialization"""
        exporter = PRDExporter(multiple_categories)

        assert exporter.categories == multiple_categories

    def test_generate_excel(self, client, category_with_features):
        """Test generating Excel workbook (no PRD Summary, only category sheets)"""
        category, features = category_with_features
        exporter = PRDExporter([category])

        excel_buffer = exporter.generate_excel()

        assert isinstance(excel_buffer, BytesIO)
        excel_buffer.seek(0)

        # Verify it's a valid workbook with category sheets (no PRD Summary)
        wb = load_workbook(excel_buffer)
        assert 'PRD Summary' not in wb.sheetnames
        assert category.name in wb.sheetnames

    def test_generate_word(self, client, category_with_features):
        """Test generating Word document"""
        category, features = category_with_features
        exporter = PRDExporter([category])

        word_buffer = exporter.generate_word()

        assert isinstance(word_buffer, BytesIO)
        word_buffer.seek(0)

        # Verify it's a valid document
        doc = Document(word_buffer)
        assert len(doc.paragraphs) > 0
        # Should contain category name
        text_content = ' '.join([p.text for p in doc.paragraphs])
        assert category.name in text_content

    def test_excel_has_correct_structure(self, client, category_with_features):
        """Test Excel file has correct structure (no PRD Summary, only category sheets)"""
        category, features = category_with_features
        exporter = PRDExporter([category])

        excel_buffer = exporter.generate_excel()
        excel_buffer.seek(0)

        wb = load_workbook(excel_buffer)

        # Should not have PRD Summary sheet (removed for backup/restore purposes)
        assert 'PRD Summary' not in wb.sheetnames

        # Check category sheet exists and has correct headers
        category_sheet = wb[category.name]
        expected_headers = ['ID', 'Title', 'Priority', 'Description', 'KPI',
                            'Customer', 'Eng. Comment',
                            'Signoff', 'Complexity', 'Release Date']

        for col, expected_header in enumerate(expected_headers, start=1):
            assert category_sheet.cell(4, col).value == expected_header

    def test_word_has_table_of_contents(self, client, category_with_features):
        """Test Word document has table of contents"""
        category, features = category_with_features
        exporter = PRDExporter([category])

        word_buffer = exporter.generate_word()
        word_buffer.seek(0)

        doc = Document(word_buffer)
        text_content = ' '.join([p.text for p in doc.paragraphs])

        # Should contain "Table of Contents"
        assert 'Table of Contents' in text_content

    def test_excel_empty_categories(self, client):
        """Test generating Excel with empty categories raises error"""
        from app.models import Category
        import pytest

        empty_category = Category(id='cat-empty', name='Empty Category', description='No features')
        exporter = PRDExporter([empty_category])

        # Should raise ValueError when no features exist
        with pytest.raises(ValueError, match="No features found to export"):
            exporter.generate_excel()
